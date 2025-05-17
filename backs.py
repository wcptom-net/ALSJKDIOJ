import os
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from openai import OpenAI
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QVBoxLayout, QPushButton,
    QLabel, QProgressBar, QComboBox, QWidget, QFileDialog
)
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QFont

# Initialize OpenAI client
API_KEY = "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"
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=API_KEY,
)

def get_ai_response(prompt):
    """Get AI-generated response for a given prompt."""
    try:
        completion = client.chat.completions.create(
            model="openai/gpt-4o",
            messages=[{"role": "user", "content": prompt}]
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"

def generate_outline_from_text(file_path):
    """Read text from a file and generate an outline using AI."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    prompt = f"請將以下文章內容轉換為大綱，使用標題和列點形式：\n\n{content}"
    return get_ai_response(prompt)

def create_ppt_from_outline(outline, output_path, color):
    """Create a PowerPoint presentation from an outline with a specific color."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "文章大綱"
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.color.rgb = color  # 設定標題顏色

    # 設定背景顏色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(
        int(color[0] * 0.9),  # 淡化顏色
        int(color[1] * 0.9),
        int(color[2] * 0.9)
    )

    for section in outline.split("\n\n"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = section.split("\n")
        if lines:
            # 設定每個章節的標題
            slide_title = slide.shapes.title
            slide_title.text = lines[0]
            slide_title.text_frame.paragraphs[0].font.color.rgb = color  # 設定標題顏色

            # 設定背景顏色
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                int(color[0] * 0.9),  # 淡化顏色
                int(color[1] * 0.9),
                int(color[2] * 0.9)
            )

            # 設定每個章節的內容
            content = "\n".join(lines[1:])  # 剩餘的行作為內容
            textbox = slide.placeholders[1]
            textbox.text = content
            for paragraph in textbox.text_frame.paragraphs:
                paragraph.font.color.rgb = color  # 設定內容文字顏色

    prs.save(output_path)

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Auto-PPT Generator")
        self.setGeometry(100, 100, 600, 400)

        # Set up the main layout
        self.central_widget = QWidget()
        self.setCentralWidget(self.central_widget)
        self.layout = QVBoxLayout(self.central_widget)

        # Add widgets
        self.label = QLabel("選擇 PPT 顏色並生成大綱")
        self.label.setFont(QFont("Arial", 16, QFont.Bold))
        self.label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.label)

        self.color_selector = QComboBox()
        self.color_selector.addItems([
            "黑色", "紅色", "藍色", "綠色", "黃色",
            "紫色", "橙色", "灰色", "粉紅色"
        ])
        self.color_selector.setStyleSheet("""
            QComboBox {
                background-color: #2b2b2b;
                color: #ffffff;
                border: 1px solid #555;
                border-radius: 5px;
                padding: 5px;
                font-size: 14px;
            }
            QComboBox::drop-down {
                border: none;
            }
            QComboBox::down-arrow {
                image: url(down_arrow.png); /* 可替換為自定義箭頭圖標 */
            }
        """)
        self.layout.addWidget(self.color_selector)

        self.folder_button = QPushButton("選擇來源資料夾")
        self.folder_button.setStyleSheet("""
            QPushButton {
                background-color: #4CAF50;
                color: white;
                border: none;
                border-radius: 5px;
                padding: 10px 20px;
                font-size: 14px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QPushButton:pressed {
                background-color: #3e8e41;
            }
        """)
        self.folder_button.clicked.connect(self.select_folder)
        self.layout.addWidget(self.folder_button)

        self.progress_bar = QProgressBar()
        self.progress_bar.setValue(0)
        self.progress_bar.setTextVisible(True)
        self.progress_bar.setStyleSheet("""
            QProgressBar {
                border: 1px solid #555;
                border-radius: 5px;
                text-align: center;
                font-size: 12px;
                background-color: #3c3c3c;
                color: #ffffff;
            }
            QProgressBar::chunk {
                background-color: #4CAF50;
                width: 20px;
            }
        """)
        self.layout.addWidget(self.progress_bar)

        self.generate_button = QPushButton("生成 PPT")
        self.generate_button.setCursor(Qt.PointingHandCursor)  # 修正 cursor 屬性
        self.generate_button.setStyleSheet("""
            QPushButton {
                background-color: #4CAF50;
                color: white;
                border: none;
                border-radius: 5px;
                padding: 10px 20px;
                font-size: 14px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QPushButton:pressed {
                background-color: #3e8e41;
            }
        """)
        self.generate_button.clicked.connect(self.generate_ppt)
        self.layout.addWidget(self.generate_button)

        # 新增結束按鈕
        self.exit_button = QPushButton("結束")
        self.exit_button.setStyleSheet("""
            QPushButton {
                background-color: #d9534f;
                color: white;
                border: none;
                border-radius: 5px;
                padding: 10px 20px;
                font-size: 14px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #c9302c;
            }
            QPushButton:pressed {
                background-color: #ac2925;
            }
        """)
        self.exit_button.clicked.connect(QApplication.quit)  # 關閉應用程式
        self.layout.addWidget(self.exit_button)

        self.status_label = QLabel("")
        self.status_label.setAlignment(Qt.AlignCenter)
        self.status_label.setStyleSheet("color: #aaaaaa; font-size: 12px;")
        self.layout.addWidget(self.status_label)

        # Apply overall CSS for Dark Mode
        self.setStyleSheet("""
            QMainWindow {
                background-color: #1e1e1e;
            }
            QLabel {
                color: #ffffff;
            }
        """)

        # Initialize folder path
        self.folder_path = ""

    def select_folder(self):
        """讓用戶選擇來源資料夾"""
        folder = QFileDialog.getExistingDirectory(self, "選擇來源資料夾", "")
        if folder:
            self.folder_path = folder
            self.status_label.setText(f"已選擇資料夾: {folder}")

    def generate_ppt(self):
        if not self.folder_path:
            self.status_label.setText("請先選擇來源資料夾！")
            return

        # 禁用按鈕並更改文字為「處理中...」
        self.generate_button.setEnabled(False)
        self.generate_button.setText("處理中...")

        text_folder = self.folder_path
        color_map = {
            0: RGBColor(0, 0, 0),       # 黑色
            1: RGBColor(255, 0, 0),     # 紅色
            2: RGBColor(0, 0, 255),     # 藍色
            3: RGBColor(0, 255, 0),     # 綠色
            4: RGBColor(255, 255, 0),   # 黃色
            5: RGBColor(128, 0, 128),   # 紫色
            6: RGBColor(255, 165, 0),   # 橙色
            7: RGBColor(128, 128, 128), # 灰色
            8: RGBColor(255, 192, 203)  # 粉紅色
        }
        selected_color = color_map[self.color_selector.currentIndex()]

        txt_files = [f for f in os.listdir(text_folder) if f.endswith(".txt")]
        total_files = len(txt_files)
        if total_files == 0:
            self.status_label.setText("沒有找到任何 .txt 檔案。")
            self.generate_button.setEnabled(True)
            self.generate_button.setText("生成 PPT")
            return

        self.progress_bar.setValue(0)
        for i, file_name in enumerate(txt_files, start=1):
            file_path = os.path.join(text_folder, file_name)
            self.status_label.setText(f"正在處理檔案: {file_name}")

            # Generate outline
            outline = generate_outline_from_text(file_path)
            if outline.startswith("Error:"):
                self.status_label.setText(f"AI 分析失敗: {outline}")
                continue

            # Create PPTX
            pptx_name = os.path.splitext(file_name)[0] + ".pptx"
            pptx_path = os.path.join(text_folder, pptx_name)
            create_ppt_from_outline(outline, pptx_path, color=selected_color)
            self.status_label.setText(f"已生成 PPTX: {pptx_name}")

            # Update progress bar
            self.progress_bar.setValue(int((i / total_files) * 100))

        self.status_label.setText("所有檔案已處理完成！")
        self.generate_button.setText("生成 PPT")
        self.generate_button.setEnabled(True)

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())